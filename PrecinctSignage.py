import pptx
from pptx import Presentation
from pptx.util import Inches
import csv
import datetime
import time
import os.path


# define inputs
file_location = "C:/Users/amalc/OneDrive/Resistor/Code/PrecinctSignage/"
preso_name = "PrecinctSignageTemplate.pptx"
csvinput_name = "PrecinctSignageTemplate.csv"

# create the presentation
prs = Presentation(file_location + preso_name)

title_slide_layout = prs.slide_layouts[0]
signage_slide_layout = prs.slide_layouts[1]

slide1 = prs.slides.add_slide(title_slide_layout)

#debug : get the slide placeholders
for shape in slide1.placeholders:
    print('%d %s' % (shape.placeholder_format.idx, shape.name))

title = slide1.shapes.title
subtitle = slide1.placeholders[1]

title.text = "Precinct Signage" 
subtitle.text= "run at " + datetime.datetime.now().strftime("%I:%M%p on %B %d, %Y")


with open(file_location + csvinput_name) as csvfile:
    reader = csv.DictReader(csvfile)

   
 
    for row in reader:
        #debug the row contents
        print(        row['organization_name'] )
        print(        row['election_name']     )
        print(        row['precinct_name']     )
        print(        row['race1_name']        )
        print(        row['race1_gfx']         )
        print(        row['race2_name']        ) 
        print(        row['race2_gfx']         )
        print(        row['race3_name']        ) 
        print(        row['race3_gfx']         )
        print(        row['race4_name']        )
        print(        row['race4_gfx']         )
        print(        row['race5_name']        )
        print(        row['race5_gfx']         )
        print(        row['race6_name']        )
        print(        row['race6_gfx']         )

        #add a slide
        slide = prs.slides.add_slide(signage_slide_layout)

        #debug : get the slide placeholders
        for shape in slide.placeholders:
               print('%d %s' % (shape.placeholder_format.idx, shape.name))

        # define the placeholders
        organization_name = slide.shapes.title #correct
        election_name = slide.placeholders[12] #correct
        precinct_name = slide.placeholders[23] #correct
        race1_name = slide.placeholders[13] #correct
        race2_name = slide.placeholders[14] #correct
        race3_name = slide.placeholders[17] #correct
        race4_name = slide.placeholders[18] #correct
        race5_name = slide.placeholders[22] #correct
        race6_name = slide.placeholders[21] #correct
        
        race1_picture = slide.placeholders[11] #correct
        race2_picture = slide.placeholders[10] #correct
        race3_picture = slide.placeholders[16] #correct
        race4_picture = slide.placeholders[15] #correct
        race5_picture = slide.placeholders[19] #correct
        race6_picture = slide.placeholders[20] #correct

        # fill the placeholders
        organization_name.text = row['organization_name']
        election_name.text = row['election_name']
        precinct_name.text = row['precinct_name']
        race1_name.text = row['race1_name']
        
        #debug
        print file_location
        print row['race1_gfx']
        race1_gfx_filename = file_location + row['race1_gfx']
        if os.path.exists(race1_gfx_filename):
             print race1_gfx_filename + "seems to be available"
        else:
             # does not exist
             print "does not seem to be available"
        race1pic = race1_picture.insert_picture(race1_gfx_filename)
        
        race2_name.text = row['race2_name']
        print row['race2_gfx']
        race2_gfx_filename = file_location + row['race2_gfx']
        if os.path.exists(race2_gfx_filename):
             print race2_gfx_filename + "seems to be available"
        else:
             # does not exist
             print race2_gfx_filename + "does not seem to be available"
        race2pic = race2_picture.insert_picture(race2_gfx_filename)
      
        race3_name.text = row['race3_name']
        print row['race3_gfx']
        race3_gfx_filename = file_location + row['race3_gfx']
        if os.path.exists(race3_gfx_filename):
             print race3_gfx_filename + "seems to be available"
        else:
             # does not exist
             print race3_gfx_filename + "does not seem to be available"
        race3pic = race3_picture.insert_picture(race3_gfx_filename)

        race4_name.text = row['race4_name']
        print row['race4_gfx']
        race4_gfx_filename = file_location + row['race4_gfx']
        if os.path.exists(race3_gfx_filename):
             print race4_gfx_filename + "seems to be available"
        else:
             # does not exist
             print race4_gfx_filename + "does not seem to be available"
        race4pic = race4_picture.insert_picture(race4_gfx_filename)

        race5_name.text = row['race5_name']
        print row['race5_gfx']
        race5_gfx_filename = file_location + row['race5_gfx']
        if os.path.exists(race5_gfx_filename):
             print race5_gfx_filename + "seems to be available"
        else:
             # does not exist
             print race5_gfx_filename + "does not seem to be available"
        race5pic = race5_picture.insert_picture(race5_gfx_filename)

        race6_name.text = row['race6_name']
        print row['race6_gfx']
        race6_gfx_filename = file_location + row['race6_gfx']
        if os.path.exists(race6_gfx_filename):
             print race6_gfx_filename + "seems to be available"
        else:
             # does not exist
             print race6_gfx_filename + "does not seem to be available"
        race6pic = race6_picture.insert_picture(race6_gfx_filename)
        
# When done save file
raw_input("Press Enter to save file and exit...")
timestr = time.strftime("%Y%m%d-%H%M%S")
filenamestr = file_location +"Precinct Signage_" + timestr +".pptx"
prs.save(filenamestr)